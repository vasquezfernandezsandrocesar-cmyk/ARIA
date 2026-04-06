"""
ARIA — Asistente Virtual Académico
Backend adaptado para Vercel (serverless, sin estado en servidor)
El historial de conversación viaja desde el frontend en cada request.
"""

import os
import json
import datetime
import tempfile
from pathlib import Path
from flask import Flask, request, jsonify
from flask_cors import CORS
import requests

# ── Lectura de documentos (opcionales) ────────────────────────────────────
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

app = Flask(__name__)
CORS(app, origins="*")

# ── Extracción de texto ───────────────────────────────────────────────────
def extract_text_from_file(filepath: str, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    text = ""
    try:
        if ext == ".pdf" and PDF_AVAILABLE:
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        elif ext == ".docx" and DOCX_AVAILABLE:
            doc = DocxDocument(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".xlsx" and XLSX_AVAILABLE:
            wb = openpyxl.load_workbook(filepath, data_only=True)
            for ws in wb.worksheets:
                text += f"\n[Hoja: {ws.title}]\n"
                for row in ws.iter_rows(values_only=True):
                    text += " | ".join([str(c) if c is not None else "" for c in row]) + "\n"
        elif ext == ".pptx" and PPTX_AVAILABLE:
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                text += f"\n[Diapositiva {i+1}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        else:
            text = f"[Formato no soportado: {filename}]"
    except Exception as e:
        text = f"[Error leyendo {filename}: {e}]"
    return text[:12000]

# ── Prompts del sistema ───────────────────────────────────────────────────
SYSTEM_BASE = """Eres ARIA, un asistente virtual académico avanzado en español (Perú).
Eres experto en psicología clínica DSM-5, redacción académica APA 7, investigación científica,
análisis crítico, productividad y educación universitaria.
Respondes siempre en español peruano, de forma clara, precisa y profesional.
Puedes conversar libremente sobre cualquier tema además de tus comandos especializados.
Fecha y hora actual en Lima, Perú: {datetime}
"""

MODES = {
    "normal":     "Responde de forma amigable, clara y directa.",
    "académico":  "Usa lenguaje formal, citas APA 7 y terminología especializada.",
    "clínico":    "Aplica DSM-5, criterios diagnósticos y lenguaje psicopatológico.",
    "creativo":   "Usa lenguaje expresivo, metáforas y narración creativa.",
}

COMMAND_PROMPTS = {
    "/noticias":      "Proporciona un resumen de las noticias más relevantes del Perú de hoy, organizado por categorías: política, economía, cultura, deportes, tecnología.",
    "/noticiasperu":  "Proporciona un resumen de las noticias más relevantes del Perú de hoy, organizado por categorías.",
    "/noticiasmundo": "Proporciona las noticias internacionales más importantes de hoy, organizadas por región: América, Europa, Asia, África.",
    "/paciente":      "Simula un caso clínico realista basado en DSM-5. Crea un paciente ficticio con: datos demográficos, motivo de consulta, historia clínica, síntomas principales. El estudiante hará la entrevista diagnóstica.",
    "/nuevopaciente": "Crea un nuevo caso clínico diferente, con un trastorno DSM-5 distinto al anterior.",
    "/feedback":      "Analiza la entrevista clínica realizada. Da retroalimentación sobre: rapport, preguntas diagnósticas, hipótesis, omisiones importantes.",
    "/salirclinico":  "Cierra la sesión clínica. Proporciona: diagnóstico correcto del caso, habilidades demostradas, áreas de mejora, calificación 1-10.",
    "/informe":       "Genera un informe académico completo en APA 7 con: portada, introducción, desarrollo, conclusiones, referencias. Si no se especificó tema, pídelo.",
    "/ensayo":        "Redacta un ensayo argumentativo universitario: introducción con tesis, 3 argumentos desarrollados, contraargumento y refutación, conclusión. APA 7.",
    "/tesis":         "Desarrolla planteamiento de tesis: título, problema, justificación, objetivos, hipótesis, variables, marco teórico inicial.",
    "/introduccion":  "Redacta introducción académica formal: contextualización, problema, importancia, objetivos, estructura del trabajo.",
    "/conclusion":    "Redacta conclusión estructurada: síntesis de hallazgos, verificación de objetivos, limitaciones, recomendaciones.",
    "/abstract":      "Genera resumen ejecutivo bilingüe (español/inglés) APA 7. Máx 250 palabras. Incluye: objetivo, método, resultados, conclusiones, palabras clave.",
    "/carta":         "Redacta carta profesional o académica formal con membrete, fecha Lima-Perú, saludo formal, cuerpo estructurado, despedida.",
    "/parafrasear":   "Reescribe el texto con 0% de similitud textual manteniendo el significado. Cambia estructura, vocabulario y orden de ideas.",
    "/resumen":       "Genera resumen académico con: idea central, ideas nucleares numeradas, conclusión sintética. Máx 30% del texto original.",
    "/mapa":          "Crea mapa conceptual en estructura de árbol ASCII: nodo central → ramas → sub-ramas → hojas.",
    "/esquema":       "Genera esquema de estudio detallado con numeración jerárquica (I, A, 1, a).",
    "/comparar":      "Crea tabla comparativa en markdown con filas de criterios y columnas de elementos.",
    "/critica":       "Análisis crítico académico: contextualización, análisis interno, valoración externa, aportes, limitaciones, juicio crítico.",
    "/argumentos":    "Presenta argumentos a favor y en contra en tabla estructurada con fuentes teóricas.",
    "/buscar":        "Realiza revisión bibliográfica estructurada con fuentes académicas, bases de datos, autores clave, tendencias.",
    "/definir":       "Definición académica multiteoría: etimología, definiciones por autores/escuelas, definición integradora.",
    "/autor":         "Biografía académica: datos biográficos, formación, obras principales, aportes teóricos, contexto histórico.",
    "/revista":       "Lista revistas científicas indexadas (Scopus, WoS, Latindex) por disciplina con ISSN y factor de impacto.",
    "/hipotesis":     "Formula hipótesis: nula, alternativa, estadística. Incluye variables y operacionalización.",
    "/metodologia":   "Diseña metodología completa: enfoque, tipo, diseño, población, muestra, instrumentos, análisis.",
    "/estadistica":   "Analiza datos estadísticos: descriptiva, pruebas de hipótesis, interpretación de resultados.",
    "/osint":         "Evalúa fuentes académicas: credibilidad, indexación, factor de impacto, sesgos, recomendación.",
    "/apa":           "Genera referencia bibliográfica completa en APA 7. Solicita tipo (libro, artículo, tesis, web) y datos necesarios.",
    "/apa_web":       "Genera cita APA 7 desde URL. Extrae autor, fecha, título, sitio web.",
    "/corregirapa":   "Corrige errores en la cita APA dada. Explica cada corrección con la regla APA 7.",
    "/bibliografia":  "Genera lista de 10+ referencias en APA 7 sobre el tema indicado.",
    "/refs":          "Genera 5 referencias rápidas en APA 7 sobre el tema.",
    "/preguntas":     "Banco de 10 preguntas: 3 conocimiento, 3 comprensión, 2 aplicación, 2 análisis (Bloom). Con respuestas.",
    "/flashcards":    "Crea 10 tarjetas de estudio pregunta-respuesta concisa para memorización activa.",
    "/plan":          "Plan de estudio personalizado: diagnóstico, objetivos, cronograma semanal, recursos, estrategias.",
    "/tecnicas":      "Técnicas de aprendizaje activo: SQ3R, Cornell, Feynman, mapas mentales, práctica espaciada.",
    "/pomodoro":      "Guía Pomodoro universitaria: protocolo 25/5, adaptaciones para exámenes, gestión de distracciones.",
    "/glosario":      "Glosario de 15 términos clave del tema con definiciones académicas.",
    "/objetivos":     "Objetivos de aprendizaje según Taxonomía de Bloom (6 niveles) con verbos específicos.",
    "/tesauro":       "Sinónimos académicos y términos relacionados para enriquecer vocabulario técnico.",
    "/traducir":      "Traduce el texto académico español↔inglés manteniendo registro formal y terminología.",
    "/corregir":      "Corrige el texto: ortografía, gramática, puntuación, coherencia, estilo académico. Explica correcciones.",
    "/imagen":        "Describe detalladamente una imagen educativa que se podría generar sobre el tema.",
}

HELP_TEXT = """
🤖 **ARIA — Asistente Virtual Académico**
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

📁 **DOCUMENTOS** — `/leer` (adjunta tu archivo)

📰 **NOTICIAS** — `/noticias` `/noticiasmundo`

🧑‍⚕️ **CLÍNICA** — `/paciente` `/nuevopaciente` `/feedback` `/salirclinico`

✍️ **REDACCIÓN** — `/informe` `/ensayo` `/tesis` `/introduccion` `/conclusion` `/abstract` `/carta` `/parafrasear`

🔬 **ANÁLISIS** — `/resumen` `/mapa` `/esquema` `/comparar` `/critica` `/argumentos`

🔍 **INVESTIGACIÓN** — `/buscar` `/definir` `/autor` `/hipotesis` `/metodologia` `/estadistica` `/osint`

📚 **APA 7** — `/apa` `/apa_web` `/corregirapa` `/bibliografia` `/refs`

🧠 **ESTUDIO** — `/preguntas` `/flashcards` `/plan` `/tecnicas` `/pomodoro` `/glosario` `/objetivos` `/tesauro`

📅 **PRODUCTIVIDAD** — `/tarea` `/tareas` `/hecha` `/recordatorio`

🔤 **HERRAMIENTAS** — `/traducir` `/corregir` `/imagen` `/hora`

⚙️ **CONFIG** — `/modo` `/stats` `/nuevo` `/ayuda`

💬 También puedes hablar libremente conmigo sobre cualquier tema.
""".strip()

# ── Comandos locales (sin IA) ─────────────────────────────────────────────
def process_local_command(cmd: str, args: str, state: dict) -> str | None:
    cmd = cmd.lower()

    if cmd == "/hora":
        now = datetime.datetime.now(datetime.timezone(datetime.timedelta(hours=-5)))
        return f"🕐 **Fecha y hora en Lima, Perú:**\n{now.strftime('%A, %d de %B de %Y')}\n🕐 {now.strftime('%H:%M:%S hrs')}"

    if cmd == "/ayuda":
        return HELP_TEXT

    if cmd == "/nuevo":
        return "__CLEAR_HISTORY__"

    if cmd == "/stats":
        return (f"📊 **Estadísticas de sesión:**\n"
                f"• Mensajes en historial: {len(state.get('history', []))}\n"
                f"• Tareas totales: {len(state.get('tasks', []))}\n"
                f"• Tareas completadas: {sum(1 for t in state.get('tasks',[]) if t.get('done'))}\n"
                f"• Modo actual: {state.get('mode','normal')}\n"
                f"• Motor IA: {state.get('provider','groq')} / {state.get('model','llama-3.3-70b-versatile')}")

    if cmd == "/tarea":
        if args:
            tasks = state.get("tasks", [])
            tasks.append({"id": len(tasks)+1, "text": args, "done": False})
            return f"__UPDATE_TASKS__{json.dumps(tasks)}__\n✅ Tarea agregada: **{args}**"
        return "⚠️ Escribe la tarea: `/tarea Leer capítulo 3`"

    if cmd == "/tareas":
        tasks = state.get("tasks", [])
        if not tasks:
            return "📋 No tienes tareas. Usa `/tarea [texto]` para agregar."
        lines = ["📋 **Lista de tareas:**\n"]
        for t in tasks:
            icon = "✅" if t.get("done") else "🔲"
            lines.append(f"{icon} {t['id']}. {t['text']}")
        return "\n".join(lines)

    if cmd == "/hecha":
        tasks = state.get("tasks", [])
        if args:
            try:
                tid = int(args.strip())
                for t in tasks:
                    if t["id"] == tid:
                        t["done"] = True
                        return f"__UPDATE_TASKS__{json.dumps(tasks)}__\n✅ Tarea #{tid} completada: **{t['text']}**"
            except ValueError:
                for t in tasks:
                    if args.lower() in t["text"].lower():
                        t["done"] = True
                        return f"__UPDATE_TASKS__{json.dumps(tasks)}__\n✅ Completada: **{t['text']}**"
        return "⚠️ Uso: `/hecha 1`"

    if cmd == "/recordatorio":
        try:
            mins = int(args.strip())
            return f"__REMINDER__{mins}__\n⏰ Recordatorio en **{mins} minuto(s)**."
        except:
            return "⚠️ Uso: `/recordatorio 25`"

    if cmd == "/modo":
        valid = ["normal", "académico", "clínico", "creativo"]
        if args.lower() in valid:
            return f"__SET_MODE__{args.lower()}__\n⚙️ Modo cambiado a: **{args.lower()}**"
        return f"⚙️ Modos: {', '.join(valid)}\nUso: `/modo académico`"

    return None

# ── Llamadas a modelos ────────────────────────────────────────────────────
def call_groq(api_key: str, model: str, messages: list, system: str) -> str:
    if not api_key:
        return "⚠️ Configura tu API Key de Groq en el panel izquierdo."
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "messages": [{"role": "system", "content": system}] + messages[-18:],
        "max_tokens": 2048,
        "temperature": 0.7,
    }
    try:
        r = requests.post("https://api.groq.com/openai/v1/chat/completions",
                          headers=headers, json=payload, timeout=55)
        data = r.json()
        if r.status_code == 200:
            return data["choices"][0]["message"]["content"]
        if r.status_code == 429:
            return "⚠️ Límite de Groq alcanzado. Espera unos segundos."
        return f"❌ Error Groq {r.status_code}: {data.get('error',{}).get('message', str(data))}"
    except requests.exceptions.Timeout:
        return "⏱️ Tiempo agotado. Intenta de nuevo."
    except Exception as e:
        return f"❌ Error: {e}"

def call_gemini(api_key: str, model: str, messages: list, system: str) -> str:
    if not api_key:
        return "⚠️ Configura tu API Key de Gemini en el panel izquierdo."
    contents = []
    for i, msg in enumerate(messages[-18:]):
        role = "user" if msg["role"] == "user" else "model"
        text = msg["content"]
        if i == 0 and role == "user":
            text = system + "\n\n" + text
        contents.append({"role": role, "parts": [{"text": text}]})
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    try:
        r = requests.post(url, json={"contents": contents,
                                      "generationConfig": {"maxOutputTokens": 2048, "temperature": 0.7}},
                          timeout=55)
        data = r.json()
        if r.status_code == 200:
            return data["candidates"][0]["content"]["parts"][0]["text"]
        if r.status_code == 429:
            return "⚠️ Límite de Gemini alcanzado. Espera."
        return f"❌ Error Gemini {r.status_code}: {data.get('error',{}).get('message', str(data))}"
    except Exception as e:
        return f"❌ Error Gemini: {e}"

def build_system(mode: str, doc_context: str) -> str:
    now = datetime.datetime.now(datetime.timezone(datetime.timedelta(hours=-5)))
    dt = now.strftime("%A, %d de %B de %Y — %H:%M hrs (Lima, Perú)")
    sys = SYSTEM_BASE.format(datetime=dt)
    sys += "\n\n" + MODES.get(mode, MODES["normal"])
    if doc_context:
        sys += f"\n\nCONTEXTO DEL DOCUMENTO CARGADO:\n{doc_context[:8000]}"
    return sys

# ── Endpoints ─────────────────────────────────────────────────────────────

@app.route("/api/chat", methods=["POST"])
def chat():
    data = request.get_json(force=True)
    message    = (data.get("message") or "").strip()
    history    = data.get("history", [])      # Lista de {role, content} desde el frontend
    state      = data.get("state", {})        # {mode, tasks, provider, model, doc_context}
    api_keys   = data.get("api_keys", {})

    mode       = state.get("mode", "normal")
    tasks      = state.get("tasks", [])
    provider   = state.get("provider", "groq")
    model_name = state.get("model", "llama-3.3-70b-versatile")
    doc_ctx    = state.get("doc_context", "")

    # Detectar comando
    cmd, args = "", ""
    if message.startswith("/"):
        parts = message.split(None, 1)
        cmd   = parts[0].lower()
        args  = parts[1] if len(parts) > 1 else ""

    # 1. Comandos locales
    local_state_update = {}
    if cmd:
        result = process_local_command(cmd, args, {**state, "tasks": tasks, "history": history})
        if result:
            # Parsear directivas de estado embebidas
            response = result
            if "__CLEAR_HISTORY__" in response:
                response = response.replace("__CLEAR_HISTORY__", "🗑️ Historial borrado. ¡Empecemos de nuevo!")
                local_state_update["clear_history"] = True
            if "__SET_MODE__" in response:
                import re
                m = re.search(r"__SET_MODE__(.+?)__", response)
                if m:
                    local_state_update["mode"] = m.group(1)
                    response = re.sub(r"__SET_MODE__.+?__\n?", "", response)
            if "__UPDATE_TASKS__" in response:
                import re
                m = re.search(r"__UPDATE_TASKS__(.+?)__", response)
                if m:
                    local_state_update["tasks"] = json.loads(m.group(1))
                    response = re.sub(r"__UPDATE_TASKS__.+?__\n?", "", response)
            if "__REMINDER__" in response:
                import re
                m = re.search(r"__REMINDER__(\d+)__", response)
                if m:
                    local_state_update["reminder_mins"] = int(m.group(1))
                    response = re.sub(r"__REMINDER__.+?__\n?", "", response)
            return jsonify({"response": response.strip(), "type": "local",
                            "state_update": local_state_update})

    # 2. Enriquecer mensaje con instrucción del comando IA
    enhanced = message
    if cmd in COMMAND_PROMPTS:
        instr = COMMAND_PROMPTS[cmd]
        enhanced = f"{instr}\n\nTema/contenido: {args}" if args else instr

    # 3. Construir historial para la IA
    full_history = history + [{"role": "user", "content": enhanced}]

    # 4. Llamar modelo
    system = build_system(mode, doc_ctx)
    api_key = api_keys.get(provider, "")

    if provider == "groq":
        response = call_groq(api_key, model_name, full_history, system)
    elif provider == "gemini":
        response = call_gemini(api_key, model_name, full_history, system)
    else:
        response = "⚠️ Proveedor no reconocido."

    return jsonify({
        "response": response,
        "type": "ai",
        "model": f"{provider}/{model_name}",
        "state_update": local_state_update,
    })


@app.route("/api/upload", methods=["POST"])
def upload():
    if "file" not in request.files:
        return jsonify({"error": "Sin archivo"}), 400
    file = request.files["file"]
    if not file.filename:
        return jsonify({"error": "Nombre vacío"}), 400

    with tempfile.NamedTemporaryFile(delete=False, suffix=Path(file.filename).suffix) as tmp:
        file.save(tmp.name)
        text = extract_text_from_file(tmp.name, file.filename)

    os.unlink(tmp.name)
    preview = text[:400] + "..." if len(text) > 400 else text
    return jsonify({
        "message": f"📄 **{file.filename}** cargado exitosamente.\n\n**Vista previa:**\n{preview}\n\nAhora puedes preguntarme sobre este documento.",
        "doc_context": text,
        "filename": file.filename,
    })


@app.route("/api/ping", methods=["GET"])
def ping():
    return jsonify({"status": "ok", "service": "ARIA"})


# Vercel necesita la app como `app`
# Para desarrollo local: python api/index.py
if __name__ == "__main__":
    app.run(debug=True, port=5000)
